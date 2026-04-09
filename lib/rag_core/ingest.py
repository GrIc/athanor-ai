import hashlib
import logging
from pathlib import Path
from typing import Optional, TYPE_CHECKING

if TYPE_CHECKING:
    from lib.rag_core.ocr import OcrProcessor

logger = logging.getLogger(__name__)


def _chunk_text(text: str, source: str, collection: str, chunk_size: int = 2000, overlap: int = 200) -> list[dict]:
    """
    Split text into fixed-size chunks with overlap.
    Returns list of chunk dicts: {"text": str, "metadata": {"source": str, "project": str, "page": int, "chunk_idx": int, "md5": str}}
    """
    if not text.strip():
        return []

    chunks = []
    start = 0
    idx = 0
    text_len = len(text)

    while start < text_len:
        end = start + chunk_size
        chunk = text[start:end]

        # try to break at last newline if not at end
        if end < text_len:
            last_nl = chunk.rfind("\n")
            if last_nl > chunk_size // 2:
                end = start + last_nl + 1
                chunk = text[start:end]

        clean_chunk = chunk.strip().replace('\x00', '').replace('\0', '')
        if clean_chunk:
            md5 = hashlib.md5(clean_chunk.encode('utf-8')).hexdigest()
            chunks.append({
                "text": clean_chunk,
                "metadata": {
                    "source": source,
                    "project": collection,
                    "page": 0, # Default to 0 unless set by caller
                    "chunk_idx": idx,
                    "md5": md5
                }
            })

        start = end - overlap
        idx += 1

    return chunks


def parse_document(file: Path, collection: str, ocr: Optional["OcrProcessor"] = None) -> list[dict]:
    """
    Parse a document into RAG chunks.
    Supported formats: .pdf, .docx, .pptx, .md, .txt, images (if ocr is provided)
    """
    if not file.exists() or not file.is_file():
        logger.warning(f"File not found or is not a file: {file}")
        return []

    suffix = file.suffix.lower()
    source_name = file.name
    all_chunks = []

    if suffix == '.pdf':
        try:
            import fitz
            doc = fitz.open(str(file))
            for page_num, page in enumerate(doc, 1):
                text = page.get_text().strip()
                if len(text) < 50 and ocr:
                    text = ocr.ocr_page(page)

                if text:
                    page_chunks = _chunk_text(text, source_name, collection)
                    # Update page number for each chunk
                    for chunk in page_chunks:
                        chunk["metadata"]["page"] = page_num
                    all_chunks.extend(page_chunks)
            doc.close()
        except Exception as e:
            logger.warning(f"Failed to parse PDF {file}: {e}")

    elif suffix == '.docx':
        try:
            from docx import Document
            doc = Document(str(file))
            text = "\n\n".join(p.text for p in doc.paragraphs if p.text.strip())
            all_chunks.extend(_chunk_text(text, source_name, collection))
        except Exception as e:
            logger.warning(f"Failed to parse DOCX {file}: {e}")

    elif suffix == '.pptx':
        try:
            from pptx import Presentation
            prs = Presentation(str(file))
            parts = []
            for i, slide in enumerate(prs.slides, 1):
                slide_texts = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        slide_texts.append(shape.text_frame.text)
                if slide_texts:
                    parts.append(f"[Slide {i}]\n" + "\n".join(slide_texts))
            text = "\n\n".join(parts)
            all_chunks.extend(_chunk_text(text, source_name, collection))
        except Exception as e:
            logger.warning(f"Failed to parse PPTX {file}: {e}")

    elif suffix in ['.md', '.txt']:
        try:
            text = file.read_text(encoding='utf-8', errors='replace')
            all_chunks.extend(_chunk_text(text, source_name, collection))
        except Exception as e:
            logger.warning(f"Failed to read text file {file}: {e}")

    elif suffix in ['.jpg', '.jpeg', '.png', '.webp']:
        if ocr:
            try:
                text = ocr.ocr_image(file)
                all_chunks.extend(_chunk_text(text, source_name, collection))
            except Exception as e:
                logger.warning(f"Failed to OCR image {file}: {e}")
        else:
             logger.warning(f"Skipping image {file} because OCR is not enabled/provided.")
    else:
        logger.warning(f"Skipping unsupported file extension {suffix} for file {file}")

    return all_chunks